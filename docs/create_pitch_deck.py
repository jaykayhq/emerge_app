#!/usr/bin/env python3
"""
Emerge Pitch Deck Generator
Creates a comprehensive investor pitch deck for Emerge - The Identity-First Habit Ecosystem
Uses Naira (₦) for all financials. Follows the 12-slide layout request.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# DESIGN SYSTEM
# ============================================================================
DEEP_VIOLET = RGBColor(0x2A, 0x1B, 0x4E)
DARK_PURPLE = RGBColor(0x1A, 0x0E, 0x33)
MID_PURPLE = RGBColor(0x4A, 0x30, 0x80)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_CORAL = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
SLIDE_BG = RGBColor(0x0F, 0x0A, 0x1E)
CARD_BG = RGBColor(0x1E, 0x14, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=ACCENT_GOLD):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_metric_card(slide, left, top, width, height, label, value, accent=ACCENT_GOLD):
    card = add_shape(slide, left, top, width, height, CARD_BG, accent)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.6),
                 value, font_size=28, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), Inches(0.5),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return card


def add_table_helper(slide, rows, cols, left, top, width, height, headers, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
            
    # Set headers
    for c_idx, header in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Calibri'
            p.font.size = Pt(13)
            p.font.color.rgb = ACCENT_GOLD
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
    # Set data
    for r_idx in range(1, rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            val = str(data[r_idx - 1][c_idx])
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = SLIDE_BG
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Calibri'
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 or val == "Yes" or val.startswith("No") else PP_ALIGN.LEFT
                if val == "Yes":
                    p.font.color.rgb = ACCENT_GREEN
                    p.font.bold = True
                elif val.startswith("No"):
                    p.font.color.rgb = ACCENT_CORAL
    return table_shape


# ============================================================================
# SLIDE 1: COVER
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, SLIDE_BG)
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_GOLD)

add_text_box(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "EMERGE", font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide1, Inches(5), Inches(3.1), Inches(3.33), ACCENT_GOLD)
add_text_box(slide1, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
             "The Identity-First Habit Ecosystem", font_size=32, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1.5), Inches(4.5), Inches(10), Inches(0.8),
             "Transforming who you are, one vote at a time.", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
             "Investor Pitch Deck  •  2026", font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
             "Powered by Flutter, Firebase & Gemini AI Behavioral Science", font_size=13, color=MID_PURPLE, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 2: PROBLEM
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, SLIDE_BG)

add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "THE PROBLEM", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide2, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Habit Apps Are Broken: The Outcome Trap", font_size=36, color=WHITE, bold=True)
add_accent_line(slide2, Inches(0.8), Inches(1.9), Inches(2), ACCENT_CORAL)

# Problem Cards
problems = [
    ("92%", "Failure Rate", "of people abandon New Year's resolutions by February due to friction."),
    ("95%", "Churn Rate", "of traditional habit tracker users quit within the first 30 days."),
    ("Outcome Trap", "Focus Mismatch", "Apps track what you do (\"lose 5kg\") rather than who you wish to become (\"I am an athlete\")."),
]

for i, (stat, title, desc) in enumerate(problems):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.4)
    add_shape(slide2, x, y, Inches(3.7), Inches(3.0), CARD_BG, ACCENT_CORAL)
    add_text_box(slide2, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.8),
                 stat, font_size=32, color=ACCENT_CORAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide2, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide2, x + Inches(0.3), y + Inches(1.7), Inches(3.1), Inches(1.1),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide2, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8),
             "The Root Cause: Most productivity apps target the superficial layer of habits—outcomes and processes.\n"
             "Emerge targets the deepest layer: Identity. True behavior change is identity change.",
             font_size=14, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 3: SOLUTION
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, SLIDE_BG)

add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "THE SOLUTION", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide3, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Emerge: The First Gamified RPG Identity Engine", font_size=36, color=WHITE, bold=True)
add_accent_line(slide3, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GREEN)

# Solution Cards
solutions = [
    ("🎭", "Identity-First Onboarding", "Users don't just add checklists. They choose a Future Self path (Stoic, Scholar, Athlete, etc.) and accumulate evidence towards that persona."),
    ("🗳️", "The Identity Vote System", "Every completed habit is structured as a \"vote\" for your character. Small, daily visual evidence shifts self-beliefs naturally."),
    ("🤖", "Gemini-Powered Coach", "An embedded AI coach that analyzes habits, context, and schedules to offer tailored stacking routines and keep users accountable."),
]

for i, (icon, title, desc) in enumerate(solutions):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.4)
    add_shape(slide3, x, y, Inches(3.7), Inches(3.8), CARD_BG, ACCENT_GREEN)
    add_text_box(slide3, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.7),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.5),
                 title, font_size=18, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + Inches(0.3), y + Inches(1.7), Inches(3.1), Inches(1.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 4: MARKET
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, SLIDE_BG)

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "MARKET OPPORTUNITY", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide4, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Massive, High-Growth Untapped Market", font_size=36, color=WHITE, bold=True)
add_accent_line(slide4, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# Top metrics
add_metric_card(slide4, Inches(0.8), Inches(2.2), Inches(3.7), Inches(1.4), "Global Habit Tracking Market", "₦17.6 Trillion", ACCENT_GOLD)
add_metric_card(slide4, Inches(4.8), Inches(2.2), Inches(3.7), Inches(1.4), "Nigeria Online Fitness (29% CAGR)", "₦209 Billion", ACCENT_CYAN)
add_metric_card(slide4, Inches(8.8), Inches(2.2), Inches(3.7), Inches(1.4), "Nigeria Wellness Projection (2033)", "₦383 Billion", ACCENT_GREEN)

# Columns
add_shape(slide4, Inches(0.8), Inches(3.9), Inches(5.7), Inches(3.0), CARD_BG, MID_PURPLE)
add_text_box(slide4, Inches(1.1), Inches(4.1), Inches(5.1), Inches(0.4), "🇳🇬 The Nigerian Digital Surge", font_size=18, color=ACCENT_GOLD, bold=True)
add_bullet_text(slide4, Inches(1.1), Inches(4.6), Inches(5.1), Inches(2.1), [
    "→ 120M+ smartphone users seeking self-improvement.",
    "→ Rising awareness around mental wellness & health.",
    "→ Digital-native Gen Z & Millennial demographic base."
], font_size=13)

add_shape(slide4, Inches(6.8), Inches(3.9), Inches(5.7), Inches(3.0), CARD_BG, MID_PURPLE)
add_text_box(slide4, Inches(7.1), Inches(4.1), Inches(5.1), Inches(0.4), "🌍 The Competitive Gap", font_size=18, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide4, Inches(7.1), Inches(4.6), Inches(5.1), Inches(2.1), [
    "→ Global apps ignore localized pricing and payment systems.",
    "→ Competitors are generic and lack cultural context/rewards.",
    "→ High-data requirements of current systems lock out users."
], font_size=13)


# ============================================================================
# SLIDE 5: PRODUCT
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, SLIDE_BG)

add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "THE PRODUCT", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide5, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Inside the Emerge RPG Habit Platform", font_size=36, color=WHITE, bold=True)
add_accent_line(slide5, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# Columns representing app visuals
add_shape(slide5, Inches(0.8), Inches(2.3), Inches(5.7), Inches(4.5), CARD_BG, MID_PURPLE)
add_text_box(slide5, Inches(1.1), Inches(2.5), Inches(5.1), Inches(0.4), "🎮 Future Self Character Studio", font_size=20, color=ACCENT_GOLD, bold=True)
add_bullet_text(slide5, Inches(1.1), Inches(3.0), Inches(5.1), Inches(3.5), [
    "→ Choose 6 distinct archetypes: Athlete, Scholar, Creator, Stoic, Explorer, Zealot.",
    "→ Dynamic visual progress: Complete habits to build skyscrapers in your City or plant trees in your Forest.",
    "→ RPG-level system: Earn XP, gold, and visual title badges with consistent execution."
], font_size=13)

add_shape(slide5, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.5), CARD_BG, MID_PURPLE)
add_text_box(slide5, Inches(7.1), Inches(2.5), Inches(5.1), Inches(0.4), "🔗 Habit Stack & Automation", font_size=20, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide5, Inches(7.1), Inches(3.0), Inches(5.1), Inches(3.5), [
    "→ Visual Habit Stack Builder: Anchors new actions onto existing daily routines.",
    "→ Gemini Coach integration: Contextual reminders, custom routines, and behavior checks.",
    "→ Mobile widgets keep habit cues prominent on user home screens."
], font_size=13)


# ============================================================================
# SLIDE 6: BUSINESS MODEL
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, SLIDE_BG)

add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "BUSINESS MODEL", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide6, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Freemium Pricing & Local Monetization", font_size=36, color=WHITE, bold=True)
add_accent_line(slide6, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# 3 pricing boxes
pricing_models = [
    ("₦2,500 / Month", "Premium Sub", "Affordable local payment, unlocking unlimited habits, heatmaps, and advanced AI logs.", ACCENT_CYAN),
    ("₦20,000 / Year", "Annual Sub", "Save 33% (equivalent to 4 months free) for long-term consistency.", ACCENT_GOLD),
    ("Challenges & Ads", "Affiliate / Ad Revenue", "Non-paying users monetize via unobtrusive ads. Affiliate challenge rewards sponsored by brands.", ACCENT_GREEN),
]

for i, (price, title, desc, color) in enumerate(pricing_models):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.3)
    add_shape(slide6, x, y, Inches(3.7), Inches(2.0), CARD_BG, color)
    add_text_box(slide6, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.4),
                 price, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x + Inches(0.2), y + Inches(0.7), Inches(3.3), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x + Inches(0.2), y + Inches(1.1), Inches(3.3), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Feature table
features_headers = ["Feature", "Free Tier", "Emerge Pro ✨"]
features_data = [
    ["Active Habits", "Max 3 Habits", "Unlimited"],
    ["Habit Stacks", "1 Morning Routine Stack", "Unlimited Stacks (Morning/Noon/Night)"],
    ["Analytics", "Current Streak", "Full Heatmaps, Trends, & Success %"],
    ["Visual Worlds", "Basic Avatar / World", "Exclusive Themes, Skins, & Auras"],
    ["AI Behavior Coach", "Basic Chat", "Advanced Coach with Full History Analysis"],
    ["Ad Experience", "Banner & Interstitial Ads", "100% Ad-Free Experience ✨"],
]
add_table_helper(slide6, 7, 3, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2), features_headers, features_data, [Inches(4.0), Inches(3.8), Inches(3.9)])


# ============================================================================
# SLIDE 7: TRACTION
# ============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, SLIDE_BG)

add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "TRACTION", font_size=14, color=ACCENT_GREEN, bold=True)
add_text_box(slide7, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Early Cohort Engagement & Proof of Value", font_size=36, color=WHITE, bold=True)
add_accent_line(slide7, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GREEN)

# Main stat card
add_shape(slide7, Inches(0.8), Inches(2.4), Inches(4.5), Inches(4.2), CARD_BG, ACCENT_GREEN)
add_text_box(slide7, Inches(1.0), Inches(2.8), Inches(4.1), Inches(1.0),
             "50+", font_size=72, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(1.0), Inches(4.0), Inches(4.1), Inches(0.4),
             "Active Beta Testers", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(1.0), Inches(4.5), Inches(4.1), Inches(0.8),
             "Providing iterative feedback on UI/UX, streak rules, and RPG progression.", font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Platform indicator inside the main card
add_shape(slide7, Inches(1.3), Inches(5.4), Inches(3.5), Inches(0.9), SLIDE_BG, MID_PURPLE)
add_text_box(slide7, Inches(1.4), Inches(5.5), Inches(3.3), Inches(0.3),
             "PLATFORMS DEPLOYED", font_size=11, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide7, Inches(1.4), Inches(5.8), Inches(3.3), Inches(0.4),
             "Android & Web (iOS Planned)", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Bullet points on right
add_shape(slide7, Inches(5.8), Inches(2.4), Inches(6.7), Inches(4.2), CARD_BG, MID_PURPLE)
add_text_box(slide7, Inches(6.1), Inches(2.7), Inches(6.1), Inches(0.4),
             "🚀 Initial Performance Indicators", font_size=20, color=ACCENT_GOLD, bold=True)
add_bullet_text(slide7, Inches(6.1), Inches(3.3), Inches(6.1), Inches(3.0), [
    "→ Over 1,000+ completed habit stacks recorded in early closed tests.",
    "→ 60% Week-1 Active User retention in initial beta cohort.",
    "→ Highly requested Private Tribe feature integrated via feedback.",
    "→ Zero system crashes reported during recent production cycles."
], font_size=14)


# ============================================================================
# SLIDE 8: COMPETITION
# ============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, SLIDE_BG)

add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "COMPETITION", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide8, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Positioning Against Global Trackers", font_size=36, color=WHITE, bold=True)
add_accent_line(slide8, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

comp_headers = ["Moat / Feature", "Emerge", "Habitica", "Fabulous", "Streaks"]
comp_data = [
    ["Identity-First Psychology", "Yes", "No (Outcome-focused)", "No (Task-focused)", "No (Goal-focused)"],
    ["Living Visual Worlds", "Yes", "Yes (Pixel Art)", "No", "No"],
    ["AI Behavior Coaching", "Yes", "No", "No", "No"],
    ["Naira/Local Pricing", "Yes", "No ($4.99/mo)", "No ($12.99/mo)", "No (Paid App Store)"],
    ["B2B Corp Challenges", "Yes", "No", "No", "No"],
    ["Offline-First Sync", "Yes", "No (Heavy latency)", "Yes", "Yes"],
]
add_table_helper(slide8, 7, 5, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.5), comp_headers, comp_data, [Inches(3.3), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)])


# ============================================================================
# SLIDE 9: GO TO MARKET
# ============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, SLIDE_BG)

add_text_box(slide9, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "GO-TO-MARKET STRATEGY", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide9, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Our Three-Phase Acquisition Flywheel", font_size=36, color=WHITE, bold=True)
add_accent_line(slide9, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# 3 horizontal cards
gtm_phases = [
    ("Phase 1: Creator Blueprints", "Partnering with top wellness, gym, and career influencers to publish premium routine plans directly on the app, drawing their audiences in.", ACCENT_CYAN),
    ("Phase 2: Social Tribes", "Viral loops driven by in-app group challenges. Members share referral codes to form private challenge tribes and win brand affiliate rewards.", ACCENT_GOLD),
    ("Phase 3: B2B Expansion", "Offering customized employee health/habit check dashboards for corporate partners, generating recurring corporate sales.", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(gtm_phases):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.4)
    add_shape(slide9, x, y, Inches(3.7), Inches(4.0), CARD_BG, color)
    add_text_box(slide9, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide9, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(2.5),
                 desc, font_size=13, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 10: TEAMS
# ============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, SLIDE_BG)

add_text_box(slide10, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "THE TEAM", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide10, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Leadership & Core Planned Hires", font_size=36, color=WHITE, bold=True)
add_accent_line(slide10, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# Left: Founder Card
add_shape(slide10, Inches(0.8), Inches(2.3), Inches(4.5), Inches(4.5), CARD_BG, ACCENT_GOLD)
add_text_box(slide10, Inches(1.1), Inches(2.6), Inches(3.9), Inches(0.4), "Sole Founder", font_size=22, color=ACCENT_GOLD, bold=True)
add_text_box(slide10, Inches(1.1), Inches(3.1), Inches(3.9), Inches(0.3), "Systems Architect & Developer", font_size=14, color=WHITE)
add_bullet_text(slide10, Inches(1.1), Inches(3.6), Inches(3.9), Inches(3.0), [
    "→ Built the complete Flutter application frontend.",
    "→ Designed the Firebase Cloud Function backend architecture.",
    "→ Deployed initial Android and Web MVP releases."
], font_size=12)

# Right: 3 planned hires vertical cards
hires = [
    ("Marketing Expert", "Drives local growth, manages creator programs, structures affiliate partnerships.", ACCENT_CYAN),
    ("Rive Animator", "Builds high-performance interactive RPG character & world animations.", ACCENT_GREEN),
    ("Business Consultant", "Opens corporate channels, secures B2B pilot licenses, guides compliance.", ACCENT_CORAL),
]

for i, (title, desc, color) in enumerate(hires):
    x = Inches(5.8)
    y = Inches(2.3 + i * 1.55)
    add_shape(slide10, x, y, Inches(6.7), Inches(1.4), CARD_BG, color)
    add_text_box(slide10, x + Inches(0.3), y + Inches(0.2), Inches(6.1), Inches(0.3),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide10, x + Inches(0.3), y + Inches(0.55), Inches(6.1), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)


# ============================================================================
# SLIDE 11: FINANCIALS
# ============================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, SLIDE_BG)

add_text_box(slide11, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "FINANCIAL PROJECTIONS", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide11, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "3-Year Conservative Growth Projections", font_size=36, color=WHITE, bold=True)
add_accent_line(slide11, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

fin_headers = ["Metric / Revenue Stream", "Year 1", "Year 2", "Year 3"]
fin_data = [
    ["Total Active Users", "50,000", "250,000", "1,000,000"],
    ["Premium Conv. Rate", "3%", "5%", "7%"],
    ["Premium Users (Pro)", "1,500", "12,500", "70,000"],
    ["Subscription Rev. (₦)", "₦45,000,000", "₦330,000,000", "₦1,680,000,000"],
    ["Ad & Affiliate Rev. (₦)", "₦15,000,000", "₦95,000,000", "₦372,000,000"],
    ["TOTAL REVENUE (₦)", "₦60,000,000", "₦425,000,000", "₦2,052,000,000"],
]
add_table_helper(slide11, 7, 4, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.5), fin_headers, fin_data, [Inches(3.7), Inches(2.6), Inches(2.7), Inches(2.7)])


# ============================================================================
# SLIDE 12: ASK
# ============================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, SLIDE_BG)

add_text_box(slide12, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "THE ASK", font_size=14, color=ACCENT_GOLD, bold=True)
add_text_box(slide12, Inches(0.8), Inches(1.0), Inches(10), Inches(0.9),
             "Pre-Seed Ask: ₦10,000,000", font_size=36, color=WHITE, bold=True)
add_accent_line(slide12, Inches(0.8), Inches(1.9), Inches(2), ACCENT_GOLD)

# Split breakdown cards (Marketing heavy)
splits = [
    ("50% - Marketing & Growth", "₦5,000,000", "Creator blueprint onboarding, local user acquisition campaigns, app store optimization (ASO).", ACCENT_CORAL),
    ("30% - Further Development", "₦3,000,000", "Contracting a Rive animator, final iOS app deployment, automated testing suite expansion.", ACCENT_CYAN),
    ("20% - B2B Expansion", "₦2,000,000", "Developing custom company challenge dashboards, structuring enterprise marketing collateral.", ACCENT_GREEN),
]

for i, (title, amount, details, color) in enumerate(splits):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.3)
    add_shape(slide12, x, y, Inches(3.7), Inches(2.2), CARD_BG, color)
    add_text_box(slide12, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide12, x + Inches(0.2), y + Inches(0.6), Inches(3.3), Inches(0.4),
                 amount, font_size=22, color=WHITE, bold=True)
    add_text_box(slide12, x + Inches(0.2), y + Inches(1.15), Inches(3.3), Inches(0.95),
                 details, font_size=11, color=LIGHT_GRAY)

# Target Milestones
add_shape(slide12, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), CARD_BG, MID_PURPLE)
add_text_box(slide12, Inches(1.1), Inches(4.9), Inches(11.1), Inches(0.3),
             "18-Month Target Milestones Following Funding Ask", font_size=14, color=ACCENT_GOLD, bold=True)

milestones = [
    ("250K", "Total Active Users"),
    ("5%", "Premium Conv."),
    ("₦425M", "Annual Revenue Target"),
    ("iOS", "Launch & App Store Presence"),
]

for i, (metric, sub) in enumerate(milestones):
    mx = Inches(1.1 + i * 2.8)
    my = Inches(5.4)
    add_shape(slide12, mx, my, Inches(2.6), Inches(1.2), SLIDE_BG, ACCENT_GOLD)
    add_text_box(slide12, mx + Inches(0.1), my + Inches(0.15), Inches(2.4), Inches(0.45),
                 metric, font_size=20, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide12, mx + Inches(0.1), my + Inches(0.65), Inches(2.4), Inches(0.45),
                 sub, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 13: THANK YOU & CONTACT INFO
# ============================================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13, SLIDE_BG)
add_shape(slide13, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT_GOLD)

add_text_box(slide13, Inches(2), Inches(1.2), Inches(9.35), Inches(1.0),
             "EMERGE", font_size=64, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide13, Inches(5), Inches(2.3), Inches(3.3), ACCENT_GOLD)

add_text_box(slide13, Inches(2), Inches(2.8), Inches(9.35), Inches(0.6),
             "Tiny changes. Remarkable results.",
             font_size=28, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_text_box(slide13, Inches(2.5), Inches(3.5), Inches(8.33), Inches(1.2),
             "Thank you for your time and consideration.\n"
             "Join us in building the next generation of habit and identity transformation.",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Contact info card
add_shape(slide13, Inches(3.5), Inches(4.9), Inches(6.3), Inches(2.1), CARD_BG, ACCENT_GOLD)
add_text_box(slide13, Inches(3.8), Inches(5.1), Inches(5.7), Inches(0.35),
             "Get In Touch", font_size=16, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide13, Inches(3.8), Inches(5.55), Inches(5.7), Inches(1.3),
             "🌐  tradeflash-l2966.web.app\n"
             "📧  joeukpai55@gmail.com\n"
             "📞  09013213794\n\n"
             "Built with ❤️ from Nigeria, for the World",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SAVE PRESENTATION
# ============================================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Emerge_Pitch_Deck.pptx")
prs.save(output_path)
print(f"[OK] Recreated pitch deck successfully saved to: {output_path}")
print(f"[INFO] Total slides: {len(prs.slides)}")
